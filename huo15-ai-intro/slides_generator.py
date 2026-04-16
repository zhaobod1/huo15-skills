#!/usr/bin/env python3
"""智能体入门到精通 PPTX 生成器
乔布斯扁平科技风 + Apple macOS 玻璃拟态卡片
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ========== 配色常量 ==========
BG_COLOR = RGBColor(0x1D, 0x1D, 0x1F)      # #1D1D1F 深灰背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # #FFFFFF 纯白
GRAY = RGBColor(0x86, 0x86, 0x8B)           # #86868B 中灰色
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF7)     # #F5F5F7 极浅灰
GLASS_FILL = RGBColor(0xFF, 0xFF, 0xFF)     # 玻璃填充（用透明度模拟）
GLASS_BORDER = RGBColor(0xFF, 0xFF, 0xFF)   # 玻璃边框

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def set_dark_background(slide):
    """设置深灰色背景"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_glass_card(slide, left, top, width, height, border_alpha=0.2):
    """添加玻璃拟态卡片"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    # 半透明白色填充（浅灰模拟玻璃效果）
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    
    # 白色细边框
    shape.line.width = Pt(0.5)
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = RGBColor(int(255 * border_alpha), int(255 * border_alpha), int(255 * border_alpha))
    
    return shape

def add_text_box(slide, left, top, width, height, text, font_size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "PingFang SC"
    p.font._element.set("{http://schemas.openxmlformats.org/drawingml/2006/main}altLang", "zh-CN")
    p.alignment = align
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ========== Slide 1: 封面 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_dark_background(slide1)
    
    # 标题
    add_text_box(slide1, Inches(0), Inches(2.5), SLIDE_W, Inches(1.2),
                 "智能体入门到精通", 72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 副标题
    add_text_box(slide1, Inches(0), Inches(3.8), SLIDE_W, Inches(0.6),
                 "成体系知识点 · 三个递进层面", 28, bold=False, color=GRAY, align=PP_ALIGN.CENTER)
    # 底部小字
    add_text_box(slide1, Inches(0), Inches(6.5), SLIDE_W, Inches(0.4),
                 "视频教程大纲", 16, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

    # ========== Slide 2: 目录概览 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide2)
    
    add_text_box(slide2, Inches(0.5), Inches(0.4), SLIDE_W, Inches(0.8),
                 "第一章 · 三层架构全景", 36, bold=True, color=WHITE)
    
    # 三个玻璃卡片
    card_w = Inches(3.8)
    card_h = Inches(2.2)
    card_top = Inches(2.2)
    gap = Inches(0.5)
    start_x = Inches(0.8)
    
    cards_data = [
        ("第一层 / 核心架构", "单体智能体四大支柱"),
        ("第二层 / 知识工程", "RAG与知识图谱"),
        ("第三层 / 协作工程化", "多智能体系统"),
    ]
    
    for i, (title, subtitle) in enumerate(cards_data):
        left = start_x + i * (card_w + gap)
        shape = add_glass_card(slide2, left, card_top, card_w, card_h)
        # 标题
        add_text_box(slide2, left + Inches(0.2), card_top + Inches(0.5), card_w - Inches(0.4), Inches(0.6),
                     title, 22, bold=True, color=WHITE)
        # 副标题
        add_text_box(slide2, left + Inches(0.2), card_top + Inches(1.2), card_w - Inches(0.4), Inches(0.6),
                     subtitle, 16, bold=False, color=GRAY)

    # ========== Slide 3: 分隔页 - 第一层 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide3)
    
    add_text_box(slide3, Inches(0), Inches(2.8), SLIDE_W, Inches(1.2),
                 "第一层", 80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide3, Inches(0), Inches(4.2), SLIDE_W, Inches(0.6),
                 "智能体核心架构", 28, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

    # ========== Slide 4: 核心架构四支柱 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide4)
    
    add_text_box(slide4, Inches(0.5), Inches(0.4), SLIDE_W, Inches(0.8),
                 "核心架构 · 四大支柱", 36, bold=True, color=WHITE)
    
    # 2x2 网格
    card_w = Inches(5.5)
    card_h = Inches(2.5)
    gap_x = Inches(1.0)
    gap_y = Inches(0.4)
    start_x = Inches(0.8)
    start_y = Inches(1.6)
    
    pillars = [
        ("规划 Planning", "任务分解 / ReAct / Chain-of-Thought"),
        ("记忆 Memory", "短期记忆 / 长期记忆 / RAG"),
        ("工具 Tool Use", "Function Calling / MCP / API集成"),
        ("大模型 LLM", "模型选型 / Fine-tuning"),
    ]
    
    for idx, (title, desc) in enumerate(pillars):
        row = idx // 2
        col = idx % 2
        left = start_x + col * (card_w + gap_x)
        top = start_y + row * (card_h + gap_y)
        
        shape = add_glass_card(slide4, left, top, card_w, card_h)
        add_text_box(slide4, left + Inches(0.3), top + Inches(0.5), card_w - Inches(0.6), Inches(0.8),
                     title, 28, bold=True, color=WHITE)
        add_text_box(slide4, left + Inches(0.3), top + Inches(1.4), card_w - Inches(0.6), Inches(0.8),
                     desc, 16, bold=False, color=GRAY)

    # ========== Slide 5: 分隔页 - 第二层 ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide5)
    
    add_text_box(slide5, Inches(0), Inches(2.8), SLIDE_W, Inches(1.2),
                 "第二层", 80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide5, Inches(0), Inches(4.2), SLIDE_W, Inches(0.6),
                 "知识工程与检索增强", 28, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

    # ========== Slide 6: 知识工程三要素 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide6)
    
    add_text_box(slide6, Inches(0.5), Inches(0.4), SLIDE_W, Inches(0.8),
                 "知识工程与检索增强", 36, bold=True, color=WHITE)
    
    # 三个卡片横向
    card_w = Inches(3.8)
    card_h = Inches(3.5)
    card_top = Inches(1.8)
    gap = Inches(0.5)
    start_x = Inches(0.8)
    
    knowledge = [
        ("RAG 检索增强生成", "向量化 / 分块策略 / 混合检索 / 重排序"),
        ("知识工程", "从文档到知识图谱，构建业务规则库"),
        ("微调 Fine-tuning", "全参数微调 / PEFT / LoRA"),
    ]
    
    for i, (title, desc) in enumerate(knowledge):
        left = start_x + i * (card_w + gap)
        shape = add_glass_card(slide6, left, card_top, card_w, card_h)
        add_text_box(slide6, left + Inches(0.2), card_top + Inches(0.5), card_w - Inches(0.4), Inches(0.8),
                     title, 22, bold=True, color=WHITE)
        add_text_box(slide6, left + Inches(0.2), card_top + Inches(1.5), card_w - Inches(0.4), Inches(1.8),
                     desc, 15, bold=False, color=GRAY)

    # ========== Slide 7: 分隔页 - 第三层 ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide7)
    
    add_text_box(slide7, Inches(0), Inches(2.8), SLIDE_W, Inches(1.2),
                 "第三层", 80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide7, Inches(0), Inches(4.2), SLIDE_W, Inches(0.6),
                 "多智能体协作与工程化", 28, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

    # ========== Slide 8: 多智能体系统 ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide8)
    
    add_text_box(slide8, Inches(0.5), Inches(0.4), SLIDE_W, Inches(0.8),
                 "多智能体协作与工程化", 36, bold=True, color=WHITE)
    
    card_w = Inches(3.8)
    card_h = Inches(3.5)
    card_top = Inches(1.8)
    gap = Inches(0.5)
    start_x = Inches(0.8)
    
    mas_data = [
        ("多智能体系统 MAS", "角色分化 / 协作模式 / AutoGPT / CAMEL框架"),
        ("协作协议", "MCP协议 / A2A协议 / AG-UI协议"),
        ("评估与稳定性", "Ragas评估 / Guardrails / 可观测性"),
    ]
    
    for i, (title, desc) in enumerate(mas_data):
        left = start_x + i * (card_w + gap)
        shape = add_glass_card(slide8, left, card_top, card_w, card_h)
        add_text_box(slide8, left + Inches(0.2), card_top + Inches(0.5), card_w - Inches(0.4), Inches(0.8),
                     title, 22, bold=True, color=WHITE)
        add_text_box(slide8, left + Inches(0.2), card_top + Inches(1.5), card_w - Inches(0.4), Inches(1.8),
                     desc, 15, bold=False, color=GRAY)

    # ========== Slide 9: 学习路径 ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide9)
    
    add_text_box(slide9, Inches(0.5), Inches(0.4), SLIDE_W, Inches(0.8),
                 "从入门到精通 · 学习路径", 36, bold=True, color=WHITE)
    
    # 时间轴 - 三个节点
    node_y = Inches(3.5)
    node_w = Inches(3.0)
    node_h = Inches(1.8)
    gap = Inches(0.8)
    start_x = Inches(1.2)
    
    timeline = [
        ("核心架构", "四大支柱构建单体智能体"),
        ("知识工程", "RAG与知识图谱增强"),
        ("协作工程化", "多智能体系统"),
    ]
    
    for i, (title, desc) in enumerate(timeline):
        left = start_x + i * (node_w + gap)
        shape = add_glass_card(slide9, left, node_y, node_w, node_h)
        add_text_box(slide9, left + Inches(0.2), node_y + Inches(0.3), node_w - Inches(0.4), Inches(0.6),
                     title, 24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide9, left + Inches(0.2), node_y + Inches(1.0), node_w - Inches(0.4), Inches(0.6),
                     desc, 14, bold=False, color=GRAY, align=PP_ALIGN.CENTER)
        
        # 箭头（用文字箭头）
        if i < 2:
            arrow_x = left + node_w + Inches(0.1)
            add_text_box(slide9, arrow_x, node_y + Inches(0.6), gap - Inches(0.2), Inches(0.6),
                         "→", 36, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ========== Slide 10: 结尾页 ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide10)
    
    add_text_box(slide10, Inches(0), Inches(2.8), SLIDE_W, Inches(1.2),
                 "开始你的智能体之旅", 60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide10, Inches(0), Inches(4.5), SLIDE_W, Inches(0.6),
                 "视频教程 · 持续更新", 24, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

    # ========== 保存 ==========
    output_path = os.path.expanduser("~/Desktop/智能体入门到精通-教程大纲.pptx")
    prs.save(output_path)
    print(f"PPTX 已生成: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
