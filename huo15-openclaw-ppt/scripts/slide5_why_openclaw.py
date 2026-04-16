#!/usr/bin/env python3
"""
合并版 PPTX — Slide 1–6（我们的公司）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_BG      = RGBColor(0x06, 0x0D, 0x1A)
C_CARD    = RGBColor(0x0D, 0x18, 0x2A)
C_TEXT    = RGBColor(0xFF, 0xFF, 0xFF)
C_SUBTEXT = RGBColor(0x88, 0x88, 0x88)
C_LIGHT   = RGBColor(0xCC, 0xCC, 0xCC)
FONT = "PingFang SC"

def text_box(slide, text, left, top, width, height,
             font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = color or C_TEXT
    return tb

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x44)
    shape.line.width = Pt(0.5)
    return shape

def add_divider(slide, left, top, width, color=None):
    ln = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(0.008)
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = color or RGBColor(0x33, 0x33, 0x44)
    ln.line.fill.background()

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════
# Slide 1 — 封面
# ════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg = s1.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = C_BG
text_box(s1, "走向具身智能", 0.8, 2.0, 11.73, 1.2, font_size=64, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
text_box(s1, "龙虾生态战略·重塑所有企业", 0.8, 3.35, 11.73, 0.7, font_size=26, color=C_TEXT, align=PP_ALIGN.CENTER)
text_box(s1, "青岛火一五信息科技有限公司  ·  2026", 0.8, 4.5, 11.73, 0.5, font_size=14, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# Slide 2 — 人工智能五个阶段
# ════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg2 = s2.background; fill2 = bg2.fill; fill2.solid(); fill2.fore_color.rgb = C_BG
text_box(s2, "人工智能的五个阶段", 0.6, 0.35, 8, 0.55, font_size=28, bold=True, color=C_TEXT)
text_box(s2, "FIVE STAGES OF ARTIFICIAL INTELLIGENCE", 0.6, 0.82, 10, 0.35, font_size=10, color=C_SUBTEXT)

stages = [
    ("1", "大数据时代", "通过大数据打造智能聊天系统", "ChatGPT", "2020–2024"),
    ("2", "推理模型时代", "不再通过大数据进行智能化，大模型拥有了推理能力", "DeepSeek", "2024"),
    ("3", "智能体时代", "人工智能不再只是拥有大模型（智力），同时具备性格/灵魂、记忆、技能、工具综合体", "龙虾 OpenClaw", "2026"),
    ("4", "具身智能时代", "智能体装到具体的机械结构上，真正意义上实现机器人时代", "宇树机器人（三年后）", "未来"),
    ("5", "镜像世界", "人工智能复刻每个人的数据行为、记忆、大脑，计算机里有一个一模一样的我们", "——", "10年后"),
]
card_w = 12.13; card_h = 0.88; gap = 0.1; start_y = 1.35
for i, (num, title, desc, rep, year) in enumerate(stages):
    y = start_y + i * (card_h + gap)
    add_card(s2, 0.6, y, card_w, card_h)
    dot = s2.shapes.add_shape(9, Inches(0.75), Inches(y + 0.3), Inches(0.28), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = C_TEXT; dot.line.fill.background()
    ntb = s2.shapes.add_textbox(Inches(0.75), Inches(y + 0.3), Inches(0.28), Inches(0.28))
    ntf = ntb.text_frame; np_ = ntf.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run(); nr.text = num; nr.font.size = Pt(11); nr.font.bold = True; nr.font.name = FONT; nr.font.color.rgb = C_BG
    text_box(s2, title, 1.5, y + 0.08, 2.5, 0.4, font_size=14, bold=True, color=C_TEXT)
    text_box(s2, desc, 1.5, y + 0.45, 7.5, 0.38, font_size=10, color=C_SUBTEXT)
    text_box(s2, "代表：" + rep, 9.1, y + 0.08, 3.3, 0.38, font_size=10, color=C_LIGHT)
    text_box(s2, year, 11.7, y + 0.45, 0.9, 0.38, font_size=10, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)
text_box(s2, "青岛火一五信息科技有限公司", 0.6, 7.1, 5, 0.3, font_size=9, color=C_SUBTEXT)
text_box(s2, "02", 12.6, 7.1, 0.5, 0.3, font_size=9, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# Slide 3 — 战略参考读物
# ════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg3 = s3.background; fill3 = bg3.fill; fill3.solid(); fill3.fore_color.rgb = C_BG
text_box(s3, "战略参考读物", 0.6, 0.35, 8, 0.55, font_size=28, bold=True, color=C_TEXT)
text_box(s3, "STRATEGIC REFERENCE BOOKS", 0.6, 0.82, 10, 0.35, font_size=10, color=C_SUBTEXT)
add_card(s3, 0.6, 1.35, 5.8, 5.4)
s3.shapes.add_picture('/tmp/book_5000_final.png', Inches(1.3), Inches(1.7), Inches(2.6), Inches(3.71))
text_box(s3, "《5000天后的世界》", 0.9, 5.55, 5.2, 0.45, font_size=14, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
text_box(s3, "凯文·凯利（Kevin Kelly）", 0.9, 5.95, 5.2, 0.35, font_size=11, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
text_box(s3, "AI、互联网、智能体时代的演进预判", 0.9, 6.3, 5.2, 0.4, font_size=10, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
add_card(s3, 6.9, 1.35, 5.8, 5.4)
s3.shapes.add_picture('/tmp/book_1000_final.jpg', Inches(7.5), Inches(1.7), Inches(3.0), Inches(3.0))
text_box(s3, "《预测之书：1000天后的世界》", 7.2, 5.0, 5.2, 0.45, font_size=14, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
text_box(s3, "罗振宇", 7.2, 5.4, 5.2, 0.35, font_size=11, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
text_box(s3, "1000天后的世界发展趋势预测", 7.2, 5.75, 5.2, 0.4, font_size=10, color=C_SUBTEXT, align=PP_ALIGN.CENTER)
text_box(s3, "青岛火一五信息科技有限公司", 0.6, 7.1, 5, 0.3, font_size=9, color=C_SUBTEXT)
text_box(s3, "03", 12.6, 7.1, 0.5, 0.3, font_size=9, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# Slide 4 — 乔布斯
# ════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg4 = s4.background; fill4 = bg4.fill; fill4.solid(); fill4.fore_color.rgb = C_BG
text_box(s4, "远见与品味", 0.6, 0.35, 8, 0.55, font_size=28, bold=True, color=C_TEXT)
text_box(s4, "VISION AND TASTE", 0.6, 0.82, 10, 0.35, font_size=10, color=C_SUBTEXT)
add_divider(s4, 0.6, 1.18, 12.13)
s4.shapes.add_picture('/tmp/steve_jobs.png', Inches(0.6), Inches(1.5), Inches(12.13), Inches(3.2))
add_card(s4, 0.6, 4.85, 12.13, 1.6)
text_box(s4, "Steve Jobs", 1.0, 5.0, 3, 0.4, font_size=16, bold=True, color=C_TEXT)
text_box(s4, "苹果公司创始人", 1.0, 5.38, 3, 0.35, font_size=11, color=C_SUBTEXT)
text_box(s4, "「找不到方向的根本原因，", 4.0, 5.0, 8.5, 0.45, font_size=22, bold=True, color=C_TEXT)
text_box(s4, "不够聪明，是没有品味。」", 4.0, 5.45, 8.5, 0.45, font_size=22, bold=True, color=C_TEXT)
text_box(s4, "—— 乔布斯", 4.0, 5.95, 8.5, 0.35, font_size=12, color=C_SUBTEXT)
text_box(s4, "青岛火一五信息科技有限公司", 0.6, 7.1, 5, 0.3, font_size=9, color=C_SUBTEXT)
text_box(s4, "04", 12.6, 7.1, 0.5, 0.3, font_size=9, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════
# Slide 5 — 为什么押宝龙虾
# ════════════════════════════════════════════════════════

# ════════════════════════════════════════════════════════
# Slide 6 — 我们的公司
# ════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg6 = s6.background; fill6 = bg6.fill; fill6.solid(); fill6.fore_color.rgb = C_BG
text_box(s6, "我们的公司", 0.6, 0.35, 8, 0.55, font_size=28, bold=True, color=C_TEXT)
text_box(s6, "COMPANIES", 0.6, 0.82, 10, 0.35, font_size=10, color=C_SUBTEXT)
add_divider(s6, 0.6, 1.18, 12.13)

# 公司一 — 青岛火一五信息科技有限公司
add_card(s6, 0.6, 1.4, 5.9, 4.5)
text_box(s6, "青岛火一五信息科技有限公司", 0.85, 1.6, 5.4, 0.45, font_size=13, bold=True, color=C_TEXT)
text_box(s6, "FIREINFO TECH", 0.85, 2.0, 5.4, 0.3, font_size=9, color=C_SUBTEXT)

company1_products = [
    ("龙虾 + 辉火云企业套件", "ERP · CRM · MES"),
    ("龙虾 + 辉火云管家", "养龙虾服务 · 数据治理"),
    ("XR-IoT 扩展现实物联网", "机器视觉 · 数字孪生"),
]
for i, (name, sub) in enumerate(company1_products):
    dot6 = s6.shapes.add_shape(9, Inches(0.85), Inches(2.55 + i * 0.85), Inches(0.18), Inches(0.18))
    dot6.fill.solid(); dot6.fill.fore_color.rgb = C_TEXT; dot6.line.fill.background()
    text_box(s6, name, 1.15, 2.48 + i * 0.85, 5, 0.35, font_size=12, bold=True, color=C_TEXT)
    text_box(s6, sub, 1.15, 2.78 + i * 0.85, 5, 0.3, font_size=10, color=C_SUBTEXT)

# 公司二 — 青岛萧伯网大科技有限公司
add_card(s6, 6.85, 1.4, 5.9, 4.5)
text_box(s6, "青岛萧伯网大科技有限公司", 7.1, 1.6, 5.4, 0.45, font_size=13, bold=True, color=C_TEXT)
text_box(s6, "XIAOBOWANG TECH", 7.1, 2.0, 5.4, 0.3, font_size=9, color=C_SUBTEXT)

text_box(s6, "逸寻智库", 7.35, 2.55, 5, 0.4, font_size=13, bold=True, color=C_TEXT)
text_box(s6, "公众号 · B站自媒体", 7.35, 2.9, 5, 0.3, font_size=10, color=C_SUBTEXT)
text_box(s6, "教育平台起步阶段", 7.35, 3.5, 5, 0.35, font_size=12, bold=True, color=C_TEXT)
text_box(s6, "布局品牌建设、IT前沿技术科普和培训", 7.35, 3.85, 5, 0.5, font_size=10, color=C_SUBTEXT)

# 个人定位
add_card(s6, 0.6, 6.1, 12.13, 0.75)
text_box(s6, "赵博  /  OPC一人公司  ·  超级个体  ·  AI工长", 1.0, 6.25, 11, 0.4, font_size=13, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

text_box(s6, "青岛火一五信息科技有限公司", 0.6, 7.1, 5, 0.3, font_size=9, color=C_SUBTEXT)
text_box(s6, "06", 12.6, 7.1, 0.5, 0.3, font_size=9, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)
s5 = prs.slides.add_slide(BLANK)
bg5 = s5.background; fill5 = bg5.fill; fill5.solid(); fill5.fore_color.rgb = C_BG
text_box(s5, "为什么押宝龙虾", 0.6, 0.35, 8, 0.55, font_size=28, bold=True, color=C_TEXT)
text_box(s5, "WHY WE BET ON OPENCLAW", 0.6, 0.82, 10, 0.35, font_size=10, color=C_SUBTEXT)
add_divider(s5, 0.6, 1.18, 12.13)

# 左卡 — 划时代现象级产品
add_card(s5, 0.6, 1.4, 5.9, 5.5)
text_box(s5, "划时代现象级产品", 0.85, 1.6, 5.4, 0.45, font_size=15, bold=True, color=C_TEXT)
add_divider(s5, 0.85, 2.1, 5.4)
text_box(s5, "龙虾作为划时代意义的现象级产品地位已经成立。", 0.85, 2.25, 5.4, 0.7, font_size=11, color=C_SUBTEXT)
text_box(s5, "它的生态、品牌、知名度不可撼动。", 0.85, 2.75, 5.4, 0.5, font_size=11, color=C_SUBTEXT)

# 繁荣度进度条
for i in range(10):
    alpha = 0.3 + i * 0.07
    v = int(255 * alpha)
    add_solid = s5.shapes.add_shape(1, Inches(0.85 + i * 0.56), Inches(3.4), Inches(0.5), Inches(0.18))
    add_solid.fill.solid(); add_solid.fill.fore_color.rgb = RGBColor(v, v, v); add_solid.line.fill.background()
text_box(s5, "社区繁荣度", 0.85, 3.7, 5.4, 0.3, font_size=9, color=C_SUBTEXT, align=PP_ALIGN.CENTER)

# 数字强调
text_box(s5, "3+", 1.2, 4.1, 1.5, 0.9, font_size=52, bold=True, color=C_TEXT)
text_box(s5, "年内优势稳固", 2.7, 4.3, 3, 0.5, font_size=13, color=C_TEXT)
text_box(s5, "就算出现更优秀的产品，大多数龙虾用户会选择坐等龙虾更新，这种情况三年内不会逆转。", 0.85, 5.1, 5.4, 0.9, font_size=10.5, color=C_SUBTEXT)

# 右卡 — 数据主权
add_card(s5, 6.85, 1.4, 5.9, 5.5)
text_box(s5, "数据主权", 7.1, 1.6, 5.4, 0.45, font_size=15, bold=True, color=C_TEXT)
add_divider(s5, 7.1, 2.1, 5.4)
text_box(s5, "企业把所有数据存到龙虾，", 7.1, 2.25, 5.4, 0.45, font_size=13, bold=True, color=C_TEXT)
text_box(s5, "然后才是小红书、抖音等APP复用一部分数据。", 7.1, 2.7, 5.4, 0.5, font_size=11, color=C_SUBTEXT)
text_box(s5, "打破以前数据割裂、数据沙箱的局面。", 7.1, 3.25, 5.4, 0.5, font_size=11, color=C_SUBTEXT)
text_box(s5, "不再让巨头瓜分我们的数据。", 7.1, 3.8, 5.4, 0.45, font_size=13, bold=True, color=C_TEXT)
text_box(s5, "龙虾 = 企业的数据金库。", 7.1, 4.35, 5.4, 0.5, font_size=16, bold=True, color=C_TEXT)
text_box(s5, "其他平台只是龙虾的附庸。", 7.1, 4.9, 5.4, 0.5, font_size=11, color=C_SUBTEXT)

text_box(s5, "青岛火一五信息科技有限公司", 0.6, 7.1, 5, 0.3, font_size=9, color=C_SUBTEXT)
text_box(s5, "05", 12.6, 7.1, 0.5, 0.3, font_size=9, color=C_SUBTEXT, align=PP_ALIGN.RIGHT)

# ─── 保存 ───────────────────────────────────────────
output = "/Users/jobzhao/.openclaw/media/outbound/合并版_走向具身智能.pptx"
prs.save(output)
print(f"✅ 合并版已生成: {output}")
