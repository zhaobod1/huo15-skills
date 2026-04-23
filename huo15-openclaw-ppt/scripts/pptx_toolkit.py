"""
pptx_toolkit.py - 按风格参数化的 PPT 绘图原语

所有函数第一个位置参数统一为 `style: Style`，便于组合复用。
上层 `create-pptx.py` 基于这些原语按 JSON 规约生成整份 deck。
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from styles import Style


def set_background(slide, style: Style):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = style.bg


def text_box(slide, style: Style, text, left, top, width, height,
             font_size=None, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size if font_size is not None else style.size_body)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = style.font
    run.font.color.rgb = color if color is not None else style.text
    return tb


def add_card(slide, style: Style, left, top, width, height, fill=None, stroke=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill is not None else style.card
    shape.line.color.rgb = stroke if stroke is not None else style.card_stroke
    shape.line.width = Pt(style.card_line_width)
    return shape


def add_divider(slide, style: Style, left, top, width, color=None, thickness=0.008):
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(thickness),
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = color if color is not None else style.divider
    ln.line.fill.background()
    return ln


def add_tag(slide, style: Style, text, left, top,
            width=1.2, height=0.32, color=None, fill=None):
    """小红书风格 hashtag 小胶囊。"""
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill if fill is not None else style.accent
    chip.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(style.size_small)
    run.font.bold = True
    run.font.name = style.font
    run.font.color.rgb = color if color is not None else style.card
    return chip


def add_accent_bar(slide, style: Style, left, top, width=0.12, height=0.45):
    """标题左侧强调色竖条（小红书常见）。"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = style.accent
    bar.line.fill.background()
    return bar


def cover_slide(prs, style: Style, title, subtitle='', footnote=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, style)

    W = style.slide_width
    H = style.slide_height

    if style.cover_decoration:
        # 小红书风格：顶部细红条 + 底部角标
        add_divider(slide, style, left=0, top=0, width=W, color=style.accent, thickness=0.08)
        add_tag(slide, style, '#火一五', left=0.5, top=H - 0.7,
                width=1.3, height=0.4)

    title_y = H * 0.32
    text_box(slide, style, title,
             left=0.6, top=title_y, width=W - 1.2, height=1.4,
             font_size=style.size_cover_title, bold=True,
             color=style.text, align=PP_ALIGN.CENTER)

    if subtitle:
        text_box(slide, style, subtitle,
                 left=0.6, top=title_y + 1.5,
                 width=W - 1.2, height=0.8,
                 font_size=style.size_cover_subtitle,
                 color=style.accent if style.name.startswith('xiaohongshu') else style.text,
                 bold=style.name.startswith('xiaohongshu'),
                 align=PP_ALIGN.CENTER)

    if footnote:
        text_box(slide, style, footnote,
                 left=0.6, top=H - 1.0, width=W - 1.2, height=0.4,
                 font_size=style.size_cover_footnote,
                 color=style.subtext, align=PP_ALIGN.CENTER)
    return slide


def content_header(slide, style: Style, title, en_subtitle=''):
    """标准内容页页首：左上角标题 + 英文副标题 + 分隔线。"""
    text_box(slide, style, title, 0.6, 0.35, 10, 0.6,
             font_size=style.size_page_title, bold=True, color=style.text)
    if style.name.startswith('xiaohongshu'):
        add_accent_bar(slide, style, left=0.35, top=0.4, width=0.12, height=0.55)
    if en_subtitle:
        sub = en_subtitle.upper() if style.upper_en_subtitle else en_subtitle
        text_box(slide, style, sub, 0.6, 0.9, 10, 0.35,
                 font_size=style.size_page_subtitle_en, color=style.subtext)
    add_divider(slide, style, 0.6, 1.22, style.slide_width - 1.2)


def page_footer(slide, style: Style, company, page_no):
    if not style.show_footer:
        return
    H = style.slide_height
    W = style.slide_width
    text_box(slide, style, company, 0.6, H - 0.4, 6, 0.3,
             font_size=style.footer_company_font_size, color=style.subtext)
    text_box(slide, style, f'{page_no:02d}', W - 1.1, H - 0.4, 0.8, 0.3,
             font_size=style.footer_company_font_size, color=style.subtext,
             align=PP_ALIGN.RIGHT)


def list_slide(prs, style: Style, title, en_subtitle, items,
               company='', page_no=1):
    """编号卡片列表页。items = [{title, desc, rep?, year?}]."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, style)
    content_header(slide, style, title, en_subtitle)

    count = len(items)
    top_margin = 1.45
    bottom_margin = 0.8
    available = style.slide_height - top_margin - bottom_margin
    gap = 0.12
    card_h = max(0.7, (available - gap * (count - 1)) / count)
    card_w = style.slide_width - 1.2

    for i, item in enumerate(items):
        y = top_margin + i * (card_h + gap)
        add_card(slide, style, 0.6, y, card_w, card_h)
        # 编号圆
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.78), Inches(y + card_h / 2 - 0.18),
            Inches(0.36), Inches(0.36),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = style.accent
        dot.line.fill.background()
        num_color = style.card if style.name.startswith('xiaohongshu') or style.name == 'jobs-dark' else style.bg
        text_box(slide, style, str(i + 1),
                 left=0.78, top=y + card_h / 2 - 0.18,
                 width=0.36, height=0.36,
                 font_size=12, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER)

        text_box(slide, style, item.get('title', ''),
                 left=1.3, top=y + 0.1, width=4.5, height=0.45,
                 font_size=style.size_card_title, bold=True, color=style.text)
        if item.get('desc'):
            text_box(slide, style, item['desc'],
                     left=1.3, top=y + 0.5, width=card_w - 4.0,
                     height=card_h - 0.55,
                     font_size=style.size_small, color=style.subtext)
        if item.get('rep'):
            text_box(slide, style, '代表：' + item['rep'],
                     left=card_w - 4.2, top=y + 0.1,
                     width=3.0, height=0.4,
                     font_size=style.size_small, color=style.light)
        if item.get('year'):
            text_box(slide, style, item['year'],
                     left=card_w - 1.0, top=y + 0.1,
                     width=1.2, height=0.4,
                     font_size=style.size_small, color=style.subtext,
                     align=PP_ALIGN.RIGHT)

    page_footer(slide, style, company, page_no)
    return slide


def quote_slide(prs, style: Style, title, en_subtitle, quote,
                author='', role='', image=None, company='', page_no=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, style)
    content_header(slide, style, title, en_subtitle)

    W = style.slide_width
    top = 1.5
    if image:
        slide.shapes.add_picture(image, Inches(0.6), Inches(top),
                                 Inches(W - 1.2), Inches(3.2))
        top += 3.3

    add_card(slide, style, 0.6, top, W - 1.2, 1.8)
    if author:
        text_box(slide, style, author, 1.0, top + 0.15, 3.5, 0.5,
                 font_size=16, bold=True, color=style.text)
    if role:
        text_box(slide, style, role, 1.0, top + 0.55, 3.5, 0.4,
                 font_size=style.size_body, color=style.subtext)
    text_box(slide, style, '「' + quote + '」',
             left=4.2, top=top + 0.15, width=W - 5.0, height=1.2,
             font_size=20, bold=True, color=style.accent, italic=False)

    page_footer(slide, style, company, page_no)
    return slide


def section_slide(prs, style: Style, big_title, sub=''):
    """大字分章页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, style)
    text_box(slide, style, big_title,
             left=0.6, top=style.slide_height * 0.35,
             width=style.slide_width - 1.2, height=1.5,
             font_size=56, bold=True, color=style.accent,
             align=PP_ALIGN.CENTER)
    if sub:
        text_box(slide, style, sub,
                 left=0.6, top=style.slide_height * 0.5,
                 width=style.slide_width - 1.2, height=0.8,
                 font_size=24, color=style.subtext, align=PP_ALIGN.CENTER)
    return slide


def end_slide(prs, style: Style, title, sub='', qrcodes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, style)
    W = style.slide_width
    H = style.slide_height
    text_box(slide, style, title,
             left=0.6, top=H * 0.3, width=W - 1.2, height=1.2,
             font_size=52, bold=True, color=style.text, align=PP_ALIGN.CENTER)
    if sub:
        text_box(slide, style, sub,
                 left=0.6, top=H * 0.45, width=W - 1.2, height=0.8,
                 font_size=28, color=style.subtext, align=PP_ALIGN.CENTER)
    if qrcodes:
        # qrcodes: [{path, label}]
        count = len(qrcodes)
        qr_w = 1.5
        total = qr_w * count + (count - 1) * 0.6
        start_x = (W - total) / 2
        y = H * 0.62
        for i, qr in enumerate(qrcodes):
            x = start_x + i * (qr_w + 0.6)
            slide.shapes.add_picture(qr['path'],
                                     Inches(x), Inches(y),
                                     Inches(qr_w), Inches(qr_w))
            if qr.get('label'):
                text_box(slide, style, qr['label'],
                         left=x - 0.2, top=y + qr_w + 0.1,
                         width=qr_w + 0.4, height=0.3,
                         font_size=style.size_small, color=style.subtext,
                         align=PP_ALIGN.CENTER)
    return slide
