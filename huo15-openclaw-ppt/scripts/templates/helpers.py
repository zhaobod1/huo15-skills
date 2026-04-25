"""
templates/helpers.py - 所有模板共用的绘图原语

不同于老的 pptx_toolkit（按 Style），这里按 StylePack 的 design tokens 绘图。
"""

from typing import Optional

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from design_system import StylePack, hex_to_rgb, resolve_weight


# ============================================================
# 画布基础
# ============================================================

def set_canvas(prs, pack: StylePack):
    """给 deck 设尺寸 —— 第一次调用生效，之后 idempotent。"""
    prs.slide_width = Inches(pack.canvas.width)
    prs.slide_height = Inches(pack.canvas.height)


# ============================================================
# 文本宽度估算 & 自动 fit
# ============================================================

def _approx_text_width_in(text: str, font_size_pt: float) -> float:
    """粗估文本在给定字号下的宽度（inch）。

    Bold display 字体偏宽，宁可高估。用偏保守的比例：
      CJK = 1.1 em（字面 + 间距）
      大写 / 数字 = 0.75 em
      小写 = 0.62 em
      标点 = 0.38 em
      空格 = 0.32 em
    """
    w_em = 0.0
    for c in text:
        cp = ord(c)
        if cp > 0x2E80:           # CJK / fullwidth
            w_em += 1.1
        elif c.isspace():
            w_em += 0.32
        elif c in '.,;:!?\'"-·—':
            w_em += 0.38
        elif c.isdigit() or c.isupper():
            w_em += 0.75
        else:
            w_em += 0.62
    # em → inch (pt/72)
    return w_em * (font_size_pt / 72.0)


def fit_font_size(text: str, width_in: float, base_size_pt: float,
                  *, min_size_pt: float = 24.0,
                  max_lines: int = 1,
                  tolerance: float = 0.88) -> float:
    """返回一个不会在 `width_in × max_lines` 溢出的字号。

    用于 hero / section 这种超大字，避免 CJK 长文本换行后撞到副标题。
    tolerance 留 8% 安全余量，避免 Impress/Keynote 实际渲染比估计略宽。
    """
    budget_in = width_in * max_lines * tolerance
    size = base_size_pt
    w = _approx_text_width_in(text, size)
    if w <= budget_in:
        return size
    # 缩到刚好能容下
    size = size * (budget_in / w)
    return max(min_size_pt, size)


def new_slide(prs, pack: StylePack):
    """新建空白幻灯片 + 填背景 + 叠装饰层（grid / dot grid / scanline / corner marks）。

    装饰叠加顺序（从下往上）：
      1. 背景（纯色 or 渐变）
      2. 网格层（grid_overlay / dot_grid）
      3. 扫描线（scanline）
      4. 四角 L 型刻度（corner_marks）
      5. 开发版本戳（dev_badge）
    """
    set_canvas(prs, pack)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    dec = pack.decoration

    # 1. 背景
    if dec.gradient_bg is not None:
        # 用全幅渐变 shape 代替纯色背景
        add_gradient_rect(
            slide, 0, 0, pack.canvas.width, pack.canvas.height,
            from_hex=dec.gradient_bg[0], to_hex=dec.gradient_bg[1],
            angle_deg=dec.gradient_bg[2],
        )
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = pack.palette.rgb('bg')

    # 2. 网格 / 点阵
    if dec.grid_overlay:
        _draw_grid_overlay(slide, pack)
    elif dec.dot_grid:
        _draw_dot_grid(slide, pack)

    # 3. 扫描线
    if dec.scanline:
        _draw_scanlines(slide, pack)

    # 4. 四角刻度
    if dec.corner_marks:
        _draw_corner_marks(slide, pack)

    # 5. 开发版本戳（在最后叠加，不被其他元素盖）
    # （dev_badge 文字内容由模板决定，这里不绘制；add_dev_badge 单独调用）

    return slide


# ============================================================
# 字距（tracking）注入 — python-pptx 官方 API 不支持，打 XML
# ============================================================

def _set_run_spacing(run, tracking_em: float, font_size_pt: int):
    """给 run 注入字距 XML。tracking_em 是 em 单位（0.02 = +2%）。"""
    if tracking_em == 0:
        return
    # OOXML 的 <a:rPr spc="N"> N 是 0.01pt 为单位
    spc_pt = tracking_em * font_size_pt  # 换算到 pt
    spc_val = int(spc_pt * 100)           # OOXML 期望 1/100 pt
    rpr = run._r.get_or_add_rPr()
    rpr.set('spc', str(spc_val))


# ============================================================
# 文本框（支持字体/字距/字重/行高）
# ============================================================

def add_text(slide, pack: StylePack, text: str,
             left: float, top: float, width: float, height: float,
             *,
             font=None, font_size=None, weight='regular',
             color_key='text_primary', color=None,
             align='left', valign='top',
             tracking=None, leading=None, uppercase=False):
    """通用文本框。

    - font         指定字体名；缺省用 display_font
    - font_size    pt；缺省用 body
    - weight       'bold' / 'semibold' / 'medium' / 'regular'
    - color_key    调色板 key（text_primary / text_secondary / accent / ...）
    - color        直接传 RGBColor 则覆盖 color_key
    - align        'left' / 'center' / 'right'
    - tracking     em 单位；正值放宽，负值收紧
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # 竖向对齐
    if valign == 'middle':
        tf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
    elif valign == 'bottom':
        tf.vertical_anchor = 4

    alignment_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    p = tf.paragraphs[0]
    p.alignment = alignment_map.get(align, PP_ALIGN.LEFT)

    # 行高
    if leading is not None:
        p.line_spacing = leading

    run = p.add_run()
    run.text = text.upper() if uppercase else text

    size_pt = font_size if font_size is not None else pack.typography.body
    run.font.size = Pt(size_pt)
    run.font.bold = resolve_weight(weight)
    run.font.name = font or pack.typography.display_font
    if color is not None:
        run.font.color.rgb = color
    else:
        run.font.color.rgb = pack.palette.rgb(color_key)

    # 字距
    if tracking is not None and tracking != 0:
        _set_run_spacing(run, tracking, size_pt)

    return tb


# ============================================================
# 形状（矩形 / 圆角矩形 / 椭圆 / 直线）
# ============================================================

def add_rect(slide, pack: StylePack,
             left, top, width, height,
             *, fill_key='bg_elevated', fill=None,
             stroke_key=None, stroke=None, stroke_width=0.0,
             radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius and radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill is not None else pack.palette.rgb(fill_key)
    if stroke is not None or stroke_key is not None:
        shape.line.color.rgb = stroke if stroke is not None else pack.palette.rgb(stroke_key)
        shape.line.width = Pt(stroke_width)
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, pack: StylePack, left, top, width, height,
             *, fill_key='accent', fill=None, stroke_width=0.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill is not None else pack.palette.rgb(fill_key)
    shape.line.fill.background()
    return shape


def add_hline(slide, pack: StylePack, left, top, width,
              *, color_key='divider', color=None, thickness=0.008):
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(thickness),
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = color if color is not None else pack.palette.rgb(color_key)
    ln.line.fill.background()
    return ln


# ============================================================
# 卡片（带描边 / 假阴影 / 圆角 — 按 pack.elevation）
# ============================================================

def add_card(slide, pack: StylePack, left, top, width, height, *, fill=None):
    """按 pack.elevation 的语言绘制卡片。"""
    elev = pack.elevation
    # 假阴影（先画）
    if elev.use_fake_shadow:
        shadow = add_rect(
            slide, pack,
            left, top + elev.shadow_offset_y, width, height,
            fill=hex_to_rgb(elev.shadow_color.rstrip('00') or '#E5E5EA'),
            radius=elev.card_radius,
        )

    # 主卡片
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if elev.card_radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill is not None else hex_to_rgb(elev.card_fill)
    if elev.card_stroke_width > 0:
        shape.line.color.rgb = hex_to_rgb(elev.card_stroke_color)
        shape.line.width = Pt(elev.card_stroke_width)
    else:
        shape.line.fill.background()
    return shape


# ============================================================
# 图片（带占位符兜底）
# ============================================================

def add_image(slide, pack: StylePack, image_path,
              left, top, width, height, *, rounded=False):
    """插图。image_path 可以是 None → 绘占位方块。"""
    if image_path:
        try:
            return slide.shapes.add_picture(
                image_path, Inches(left), Inches(top),
                Inches(width), Inches(height),
            )
        except Exception:
            pass
    # 占位方块
    radius = pack.elevation.card_radius if rounded else 0
    return add_rect(
        slide, pack, left, top, width, height,
        fill_key='bg_subtle', radius=radius,
    )


# ============================================================
# 标题块（页面标题 + 英文副标题 + accent bar）
# ============================================================

def add_page_header(slide, pack: StylePack, title: str, en_sub: str = ''):
    """标准页头：左上角标题 + 英文副标题 + 可选分隔线。"""
    dec = pack.decoration
    t = pack.typography
    sp = pack.spacing
    W = pack.canvas.width

    y = 0.5

    # 英文副标题（如果是 'above' 位置）
    if en_sub and dec.page_en_sub_position == 'above':
        sub = en_sub.upper() if t.uppercase_en_sub else en_sub
        add_text(slide, pack, sub,
                 left=sp.margin_x, top=y, width=W - 2 * sp.margin_x, height=0.35,
                 font=pack.typography.body_font,
                 font_size=t.page_sub, weight='medium',
                 color_key='accent', align=dec.page_title_align,
                 tracking=0.15 if t.uppercase_en_sub else 0.02)
        y += 0.45

    # accent bar
    if dec.page_accent_bar:
        add_rect(slide, pack,
                 left=sp.margin_x, top=y + 0.05, width=0.08, height=0.55,
                 fill_key='accent')
        title_left = sp.margin_x + 0.25
    else:
        title_left = sp.margin_x

    # 主标题
    add_text(slide, pack, title,
             left=title_left, top=y,
             width=W - title_left - sp.margin_x, height=0.8,
             font_size=t.page_title, weight=t.page_weight,
             color_key='text_primary', align=dec.page_title_align,
             tracking=t.page_tracking)
    y += 0.75

    # 英文副标题（如果是 'under' 位置）
    if en_sub and dec.page_en_sub_position == 'under':
        sub = en_sub.upper() if t.uppercase_en_sub else en_sub
        add_text(slide, pack, sub,
                 left=sp.margin_x, top=y, width=W - 2 * sp.margin_x, height=0.35,
                 font=pack.typography.body_font,
                 font_size=t.page_sub, weight='regular',
                 color_key='text_tertiary', align=dec.page_title_align,
                 tracking=0.12 if t.uppercase_en_sub else 0.01)
        y += 0.4

    # 底部细分隔线（有些风格会画，有些不画）
    if pack.name.startswith('apple') or pack.name == 'jobs-dark':
        pass  # Apple 不画横线
    else:
        add_hline(slide, pack, sp.margin_x, y + 0.1, W - 2 * sp.margin_x)

    return y + 0.25


# ============================================================
# 页脚
# ============================================================

def add_page_footer(slide, pack: StylePack, company: str, page_no: int):
    if not pack.show_footer:
        return
    t = pack.typography
    sp = pack.spacing
    W = pack.canvas.width
    H = pack.canvas.height

    add_text(slide, pack, company,
             left=sp.margin_x, top=H - 0.42, width=6, height=0.3,
             font=t.body_font, font_size=t.page_number,
             color_key='text_tertiary')
    add_text(slide, pack, f'{page_no:02d}',
             left=W - sp.margin_x - 0.8, top=H - 0.42, width=0.8, height=0.3,
             font=t.body_font, font_size=t.page_number,
             color_key='text_tertiary', align='right')


# ============================================================
# v3.1 科技风装饰原语
# ============================================================

_NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _nsmap():
    return {'a': _NS_A}


def _hex_to_aval(hex_str: str) -> str:
    """'#RRGGBB' → 'RRGGBB'（OOXML srgbClr 的 val 去掉 #）。"""
    return hex_str.lstrip('#').upper()[:6]


def add_gradient_rect(slide, left, top, width, height,
                      *, from_hex: str, to_hex: str, angle_deg: int = 90):
    """全幅或局部渐变矩形（无描边）。

    OOXML 的 gradFill 线性角度单位是 1/60000 度，0° = 水平左→右。
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.line.fill.background()

    spPr_elem = shape._element.find(qn('p:spPr'))
    if spPr_elem is None:
        return shape

    # 清掉现有 fill
    for tag in ('a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:noFill'):
        old = spPr_elem.find(qn(tag))
        if old is not None:
            spPr_elem.remove(old)

    # 注入 gradFill
    angle_units = int(angle_deg * 60000)
    grad_xml = f'''<a:gradFill xmlns:a="{_NS_A}" flip="none" rotWithShape="1">
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="{_hex_to_aval(from_hex)}"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="{_hex_to_aval(to_hex)}"/></a:gs>
      </a:gsLst>
      <a:lin ang="{angle_units}" scaled="0"/>
      <a:tileRect/>
    </a:gradFill>'''
    grad_elem = etree.fromstring(grad_xml)
    # 插到 ln 之前（fill 必须在 ln 前）
    ln_elem = spPr_elem.find(qn('a:ln'))
    if ln_elem is not None:
        ln_idx = list(spPr_elem).index(ln_elem)
        spPr_elem.insert(ln_idx, grad_elem)
    else:
        spPr_elem.append(grad_elem)
    return shape


def apply_text_gradient(run, from_hex: str, to_hex: str, angle_deg: int = 0):
    """给 run.font 注入 gradFill 代替 solidFill。

    常用于 hero 大字做电青→电紫渐变。
    """
    rpr = run._r.get_or_add_rPr()
    # 移除现有 solidFill（python-pptx 会先设这个）
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill'):
        old = rpr.find(qn(tag))
        if old is not None:
            rpr.remove(old)

    angle_units = int(angle_deg * 60000)
    grad_xml = f'''<a:gradFill xmlns:a="{_NS_A}" flip="none" rotWithShape="1">
      <a:gsLst>
        <a:gs pos="0"><a:srgbClr val="{_hex_to_aval(from_hex)}"/></a:gs>
        <a:gs pos="100000"><a:srgbClr val="{_hex_to_aval(to_hex)}"/></a:gs>
      </a:gsLst>
      <a:lin ang="{angle_units}" scaled="0"/>
      <a:tileRect/>
    </a:gradFill>'''
    grad_elem = etree.fromstring(grad_xml)
    # gradFill 通常插在 rPr 的子元素列表最前（solidFill 原本位置）
    rpr.insert(0, grad_elem)


def _set_shape_alpha(shape, alpha_pct: float):
    """给 shape 的 solid fill 加 alpha（0~1）。"""
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    solid = spPr.find(qn('a:solidFill'))
    if solid is None:
        return
    clr = solid.find(qn('a:srgbClr'))
    if clr is None:
        return
    old = clr.find(qn('a:alpha'))
    if old is not None:
        clr.remove(old)
    alpha_val = int(max(0, min(1, alpha_pct)) * 100000)
    alpha = etree.SubElement(clr, qn('a:alpha'))
    alpha.set('val', str(alpha_val))


def _draw_grid_overlay(slide, pack: StylePack):
    """细线网格（横 + 竖）。"""
    dec = pack.decoration
    W, H = pack.canvas.width, pack.canvas.height
    spacing = max(dec.grid_spacing, 0.1)
    thickness = max(dec.grid_thickness, 0.003)
    color = hex_to_rgb(dec.grid_color)

    # 竖线
    x = spacing
    while x < W:
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(0),
            Inches(thickness), Inches(H),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = color
        ln.line.fill.background()
        x += spacing
    # 横线
    y = spacing
    while y < H:
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(y),
            Inches(W), Inches(thickness),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = color
        ln.line.fill.background()
        y += spacing


def _draw_dot_grid(slide, pack: StylePack):
    """圆点网格 —— 极简科技感。"""
    dec = pack.decoration
    W, H = pack.canvas.width, pack.canvas.height
    spacing = max(dec.dot_spacing, 0.15)
    size = max(dec.dot_size, 0.02)
    color = hex_to_rgb(dec.dot_color)

    y = spacing
    while y < H:
        x = spacing
        while x < W:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - size / 2), Inches(y - size / 2),
                Inches(size), Inches(size),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            x += spacing
        y += spacing


def _draw_scanlines(slide, pack: StylePack):
    """水平扫描线 —— retro-tech / CRT 感。每 0.08in 一条。"""
    dec = pack.decoration
    W, H = pack.canvas.width, pack.canvas.height
    color = hex_to_rgb(dec.scanline_color)
    y = 0.08
    while y < H:
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(y),
            Inches(W), Inches(0.004),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = color
        ln.line.fill.background()
        y += 0.08


def _draw_corner_marks(slide, pack: StylePack):
    """四角 L 型刻度（取景框感）。"""
    dec = pack.decoration
    W, H = pack.canvas.width, pack.canvas.height
    size = dec.corner_size
    thick = dec.corner_thickness
    color = pack.palette.rgb('accent')
    margin = 0.5  # 距离边缘

    def _mark(anchor_x, anchor_y, dir_x, dir_y):
        """在 (anchor_x, anchor_y) 处画一个 L，朝 (dir_x, dir_y) 方向延伸。"""
        # 横线
        hx = anchor_x if dir_x > 0 else anchor_x - size
        hy = anchor_y if dir_y > 0 else anchor_y - thick
        h = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(hx), Inches(hy),
            Inches(size), Inches(thick),
        )
        h.fill.solid()
        h.fill.fore_color.rgb = color
        h.line.fill.background()
        # 竖线
        vx = anchor_x if dir_x > 0 else anchor_x - thick
        vy = anchor_y if dir_y > 0 else anchor_y - size
        v = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(vx), Inches(vy),
            Inches(thick), Inches(size),
        )
        v.fill.solid()
        v.fill.fore_color.rgb = color
        v.line.fill.background()

    _mark(margin, margin, +1, +1)                # 左上
    _mark(W - margin, margin, -1, +1)            # 右上
    _mark(margin, H - margin, +1, -1)            # 左下
    _mark(W - margin, H - margin, -1, -1)        # 右下


def add_glow_halo(slide, pack: StylePack, cx: float, cy: float,
                  radius: float, *, color_key: str = 'accent',
                  layers: int = 4, strength: float = 0.6):
    """在 (cx, cy) 周围画多层半透明椭圆模拟辉光。

    layers 层，每层半径递增、alpha 递减。
    """
    color = pack.palette.rgb(color_key)
    for i in range(layers, 0, -1):
        r = radius * (0.4 + 0.2 * i)
        alpha = strength * (0.08 * i / layers)
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - r), Inches(cy - r * 0.6),  # 椭圆略扁
            Inches(r * 2), Inches(r * 1.2),
        )
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.fill.background()
        _set_shape_alpha(oval, alpha)


def add_dev_badge(slide, pack: StylePack, text: str, *, position: str = 'bottom-left'):
    """左下等宽字体 badge，比如 `BUILD · 2026.4.24` / `v3.1.0`。"""
    dec = pack.decoration
    sp = pack.spacing
    W, H = pack.canvas.width, pack.canvas.height

    font = dec.mono_font or dec.mono_fallbacks[0]

    if position == 'bottom-left':
        left, top = sp.margin_x, H - 0.55
        align = 'left'
    elif position == 'bottom-right':
        left, top = W - sp.margin_x - 4, H - 0.55
        align = 'right'
    elif position == 'top-right':
        left, top = W - sp.margin_x - 4, 0.35
        align = 'right'
    else:  # top-left
        left, top = sp.margin_x, 0.35
        align = 'left'

    add_text(slide, pack, text,
             left=left, top=top, width=4, height=0.3,
             font=font, font_size=pack.typography.page_number,
             weight='regular',
             color_key='text_tertiary',
             align=align,
             tracking=0.15, uppercase=True)


def add_mono_text(slide, pack: StylePack, text: str,
                  left, top, width, height,
                  *, font_size=None, color_key='text_tertiary',
                  align='left', weight='regular',
                  tracking=0.05, uppercase=False):
    """等宽字体文本（caption / metadata 用）。"""
    dec = pack.decoration
    font = dec.mono_font or dec.mono_fallbacks[0]
    return add_text(slide, pack, text,
                    left=left, top=top, width=width, height=height,
                    font=font, font_size=font_size or pack.typography.caption,
                    weight=weight, color_key=color_key, align=align,
                    tracking=tracking, uppercase=uppercase)


def add_vertical_hairline(slide, pack: StylePack, x, top, height,
                          *, color_key='border', color=None, thickness=0.006):
    """细竖线（分栏装饰用）。"""
    ln = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(top),
        Inches(thickness), Inches(height),
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = color if color is not None else pack.palette.rgb(color_key)
    ln.line.fill.background()
    return ln


def format_dev_badge(pack: StylePack, template: Optional[str] = None,
                     year: Optional[str] = None, build: Optional[str] = None) -> str:
    """根据 decoration.dev_badge_template 插值生成 badge 文本。

    占位：{year} {date} {build} {n}
    """
    tmpl = template or pack.decoration.dev_badge_template
    import datetime
    today = datetime.date.today()
    return (tmpl
            .replace('{year}', year or str(today.year))
            .replace('{date}', today.strftime('%Y.%-m.%-d')
                    if hasattr(today, 'strftime') else '2026.4.24')
            .replace('{build}', build or '0001')
            .replace('{n}', build or '0001'))


# ============================================================
# v3.2 新增：液态玻璃 / 朱砂印章 / 飞白笔触 / 几何装饰
# ============================================================

def add_color_orb(slide, cx: float, cy: float, radius: float,
                  *, color: str, alpha: float = 0.55, layers: int = 6):
    """彩色光球 —— Apple Liquid Glass 招牌装饰。

    在 (cx, cy) 处叠多层半透明椭圆，从中心向外 alpha 递减，模拟高斯模糊后的彩色光晕。
    color: '#RRGGBB'
    layers: 5~7 层最像
    """
    from pptx.dml.color import RGBColor as _RGB
    from pptx.util import Inches as _In
    from pptx.enum.shapes import MSO_SHAPE as _SH

    rgb = hex_to_rgb(color)
    for i in range(layers, 0, -1):
        r = radius * (0.35 + 0.18 * i)
        a = alpha * ((i / layers) ** 1.6) * 0.35
        orb = slide.shapes.add_shape(
            _SH.OVAL,
            _In(cx - r), _In(cy - r),
            _In(r * 2), _In(r * 2),
        )
        orb.fill.solid()
        orb.fill.fore_color.rgb = rgb
        orb.line.fill.background()
        _set_shape_alpha(orb, a)


def add_orb_cluster(slide, pack: StylePack, *,
                    palette_hex: Optional[list] = None,
                    seed: int = 7, count: int = 6):
    """在画布中铺一组分布式彩色光球（Liquid Glass / WWDC 主视觉）。

    palette_hex: ['#0A84FF', '#BF5AF2', ...] —— 默认走 Apple system colors。
    seed: 随机种子，固定可复现。
    """
    import random
    if palette_hex is None:
        palette_hex = ['#0A84FF', '#BF5AF2', '#FF375F', '#FF9F0A',
                       '#30D158', '#64D2FF', '#5E5CE6']
    rng = random.Random(seed)
    W, H = pack.canvas.width, pack.canvas.height
    for _ in range(count):
        cx = rng.uniform(W * 0.05, W * 0.95)
        cy = rng.uniform(H * 0.05, H * 0.95)
        radius = rng.uniform(1.4, 2.6)
        color = rng.choice(palette_hex)
        alpha = rng.uniform(0.4, 0.7)
        add_color_orb(slide, cx, cy, radius,
                      color=color, alpha=alpha, layers=6)


def add_glass_card(slide, left, top, width, height,
                   *, fill_hex: str = '#FFFFFF', alpha: float = 0.55,
                   stroke_hex: str = '#FFFFFF', stroke_width_pt: float = 0.5,
                   radius: float = 0.4):
    """半透磨砂玻璃卡 —— Liquid Glass 招牌。

    用 ROUNDED_RECTANGLE + alpha + 白色细边模拟磨砂质感。
    """
    from pptx.util import Inches as _In, Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE as _SH

    card = slide.shapes.add_shape(
        _SH.ROUNDED_RECTANGLE,
        _In(left), _In(top),
        _In(width), _In(height),
    )
    # 圆角调整（python-pptx 默认 adj 太小）
    try:
        card.adjustments[0] = min(0.5, radius / max(width, height))
    except Exception:
        pass
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    _set_shape_alpha(card, alpha)
    card.line.color.rgb = hex_to_rgb(stroke_hex)
    card.line.width = _Pt(stroke_width_pt)
    return card


def add_seal_stamp(slide, pack: StylePack, text: str,
                   *, left: float, top: float, size: float = 0.7,
                   color_hex: Optional[str] = None,
                   variant: str = 'square'):
    """朱砂方印 —— 原研哉 / 水墨 / 国风 共用装饰。

    variant: 'square' 阳刻方印 / 'circle' 阴刻圆印
    text: 1-4 个字（最佳）
    """
    from pptx.util import Inches as _In, Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE as _SH
    from pptx.enum.text import PP_ALIGN as _PA

    color_hex = color_hex or pack.palette.accent
    rgb = hex_to_rgb(color_hex)

    shape_kind = _SH.RECTANGLE if variant == 'square' else _SH.OVAL
    seal = slide.shapes.add_shape(
        shape_kind,
        _In(left), _In(top), _In(size), _In(size),
    )
    seal.fill.solid()
    seal.fill.fore_color.rgb = rgb
    seal.line.color.rgb = rgb
    seal.line.width = _Pt(1.5)

    # 内嵌白字
    tf = seal.text_frame
    tf.margin_left = _In(0.04)
    tf.margin_right = _In(0.04)
    tf.margin_top = _In(0.04)
    tf.margin_bottom = _In(0.04)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = _PA.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'STKaiti'
    # 字数自适应
    n = len([c for c in text if c.strip()])
    pt = 22 if n <= 1 else 18 if n == 2 else 14 if n == 3 else 11
    run.font.size = _Pt(pt)
    run.font.bold = True
    from pptx.dml.color import RGBColor as _RGB
    run.font.color.rgb = _RGB(0xFF, 0xFF, 0xF8)
    return seal


def add_brushstroke_band(slide, left: float, top: float,
                         width: float, height: float,
                         *, color_hex: str = '#1A1A1A',
                         alpha: float = 0.85,
                         tilt_deg: float = 0):
    """飞白笔触横扫 —— 水墨风招牌。

    用一个矩形 + 极小高度 + alpha 模拟书法横扫的笔触。
    多次叠加（recommended: 调用方叠 3-5 道，错位 + 不同 alpha）效果更像。
    """
    from pptx.util import Inches as _In
    from pptx.enum.shapes import MSO_SHAPE as _SH

    band = slide.shapes.add_shape(
        _SH.RECTANGLE,
        _In(left), _In(top),
        _In(width), _In(height),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = hex_to_rgb(color_hex)
    _set_shape_alpha(band, alpha)
    band.line.fill.background()
    if tilt_deg:
        try:
            band.rotation = tilt_deg
        except Exception:
            pass
    return band


def add_brushstroke_cluster(slide, pack: StylePack,
                            cx: float, cy: float,
                            *, length: float = 4.0, color_key: str = 'text_primary',
                            count: int = 5):
    """一组飞白笔触 —— 调用 add_brushstroke_band 叠加 5 道。"""
    import random
    color = getattr(pack.palette, color_key)
    rng = random.Random(int(cx * 1000) + int(cy * 1000))
    for i in range(count):
        offset_y = (i - count // 2) * 0.045 + rng.uniform(-0.02, 0.02)
        offset_x = rng.uniform(-0.15, 0.15)
        h = rng.uniform(0.025, 0.055)
        a = 0.35 + rng.uniform(0.0, 0.45)
        add_brushstroke_band(
            slide,
            left=cx - length / 2 + offset_x,
            top=cy + offset_y,
            width=length + rng.uniform(-0.4, 0.4),
            height=h,
            color_hex=color,
            alpha=a,
        )


def add_paint_stroke(slide, left: float, top: float,
                     width: float, height: float,
                     *, color_hex: str = '#FFC107',
                     alpha: float = 0.92,
                     tilt_deg: float = -8):
    """单道油画粗笔触 —— 梵高风。"""
    from pptx.util import Inches as _In
    from pptx.enum.shapes import MSO_SHAPE as _SH

    stroke = slide.shapes.add_shape(
        _SH.ROUNDED_RECTANGLE,
        _In(left), _In(top),
        _In(width), _In(height),
    )
    try:
        stroke.adjustments[0] = 0.5
    except Exception:
        pass
    stroke.fill.solid()
    stroke.fill.fore_color.rgb = hex_to_rgb(color_hex)
    _set_shape_alpha(stroke, alpha)
    stroke.line.fill.background()
    if tilt_deg:
        try:
            stroke.rotation = tilt_deg
        except Exception:
            pass
    return stroke


def add_paint_stroke_cluster(slide, pack: StylePack,
                             cx: float, cy: float,
                             *, palette_keys: Optional[list] = None,
                             count: int = 8, span: float = 4.0):
    """一组油画笔触叠加 —— 梵高星夜感。"""
    import random
    palette_keys = palette_keys or ['accent', 'accent_soft', 'text_primary']
    rng = random.Random(int(cx * 100) + int(cy * 100))
    for _ in range(count):
        color = getattr(pack.palette, rng.choice(palette_keys))
        w = rng.uniform(0.6, 1.4)
        h = rng.uniform(0.06, 0.12)
        x = cx + rng.uniform(-span / 2, span / 2)
        y = cy + rng.uniform(-1.5, 1.5)
        tilt = rng.uniform(-30, 30)
        a = rng.uniform(0.6, 0.92)
        add_paint_stroke(slide, x, y, w, h,
                         color_hex=color, alpha=a, tilt_deg=tilt)


def add_geometric_decoration(slide, pack: StylePack,
                             *, mode: str = 'memphis', seed: int = 13):
    """孟菲斯/包豪斯风的彩色几何装饰 —— 圆 / 三角 / 菱形 / 之字纹。

    mode: 'memphis' / 'bauhaus' / 'minimal'
    """
    import random
    from pptx.util import Inches as _In, Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE as _SH

    rng = random.Random(seed)
    W, H = pack.canvas.width, pack.canvas.height
    palette = [pack.palette.accent, pack.palette.accent_soft,
               pack.palette.text_primary]

    shapes_pool = [_SH.OVAL, _SH.ISOCELES_TRIANGLE, _SH.DIAMOND,
                   _SH.RIGHT_TRIANGLE, _SH.PENTAGON]
    count = 5 if mode == 'memphis' else 3 if mode == 'bauhaus' else 2

    for _ in range(count):
        kind = rng.choice(shapes_pool)
        size = rng.uniform(0.6, 1.4)
        x = rng.uniform(0.5, W - 1.5)
        y = rng.uniform(0.5, H - 1.5)
        s = slide.shapes.add_shape(
            kind, _In(x), _In(y), _In(size), _In(size),
        )
        s.fill.solid()
        s.fill.fore_color.rgb = hex_to_rgb(rng.choice(palette))
        if mode == 'memphis':
            s.line.color.rgb = hex_to_rgb(pack.palette.text_primary)
            s.line.width = _Pt(2.0)
        elif mode == 'bauhaus':
            s.line.fill.background()
        else:
            s.line.color.rgb = hex_to_rgb(pack.palette.border)
            s.line.width = _Pt(0.5)
        try:
            s.rotation = rng.uniform(0, 60)
        except Exception:
            pass


def add_chinese_pattern_border(slide, pack: StylePack,
                               *, color_hex: Optional[str] = None,
                               thickness: float = 0.025,
                               margin: float = 0.35):
    """国风万字纹/双线金边框 —— 沿画布四周。

    简化版：双线金边，外粗内细。
    """
    from pptx.util import Inches as _In
    from pptx.enum.shapes import MSO_SHAPE as _SH

    color_hex = color_hex or '#FFB61E'  # 藤黄默认
    rgb = hex_to_rgb(color_hex)
    W, H = pack.canvas.width, pack.canvas.height

    # 四条外粗边
    for x, y, w, h in [
        (margin, margin, W - 2 * margin, thickness),                 # 上
        (margin, H - margin - thickness, W - 2 * margin, thickness), # 下
        (margin, margin, thickness, H - 2 * margin),                 # 左
        (W - margin - thickness, margin, thickness, H - 2 * margin), # 右
    ]:
        e = slide.shapes.add_shape(
            _SH.RECTANGLE, _In(x), _In(y), _In(w), _In(h),
        )
        e.fill.solid()
        e.fill.fore_color.rgb = rgb
        e.line.fill.background()

    # 内细边 0.06" 内偏移
    inner = margin + 0.08
    inner_thick = thickness * 0.4
    for x, y, w, h in [
        (inner, inner, W - 2 * inner, inner_thick),
        (inner, H - inner - inner_thick, W - 2 * inner, inner_thick),
        (inner, inner, inner_thick, H - 2 * inner),
        (W - inner - inner_thick, inner, inner_thick, H - 2 * inner),
    ]:
        e = slide.shapes.add_shape(
            _SH.RECTANGLE, _In(x), _In(y), _In(w), _In(h),
        )
        e.fill.solid()
        e.fill.fore_color.rgb = rgb
        e.line.fill.background()
        _set_shape_alpha(e, 0.6)


def add_offset_shadow_block(slide, pack: StylePack,
                            left: float, top: float,
                            width: float, height: float,
                            *, fill_hex: str = '#FFFFFF',
                            shadow_hex: str = '#1A1A1A',
                            offset: float = 0.08,
                            stroke_pt: float = 2.0):
    """孟菲斯/Y2K 偏移投影块 —— 黑色实心投影 + 白底彩边。"""
    from pptx.util import Inches as _In, Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE as _SH

    # 投影
    sh = slide.shapes.add_shape(
        _SH.RECTANGLE,
        _In(left + offset), _In(top + offset),
        _In(width), _In(height),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = hex_to_rgb(shadow_hex)
    sh.line.fill.background()

    # 主块
    blk = slide.shapes.add_shape(
        _SH.RECTANGLE,
        _In(left), _In(top), _In(width), _In(height),
    )
    blk.fill.solid()
    blk.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    blk.line.color.rgb = hex_to_rgb(shadow_hex)
    blk.line.width = _Pt(stroke_pt)
    return blk


def add_golden_ratio_guide(slide, pack: StylePack,
                           *, color_hex: Optional[str] = None,
                           thickness: float = 0.005,
                           alpha: float = 0.25):
    """达芬奇/文艺复兴风：黄金分割辅助网格（0.382 / 0.618 双向）。"""
    from pptx.util import Inches as _In
    from pptx.enum.shapes import MSO_SHAPE as _SH

    color_hex = color_hex or pack.palette.text_muted
    W, H = pack.canvas.width, pack.canvas.height
    rgb = hex_to_rgb(color_hex)

    # 垂直黄金线
    for ratio in (0.382, 0.618):
        ln = slide.shapes.add_shape(
            _SH.RECTANGLE,
            _In(W * ratio), _In(0),
            _In(thickness), _In(H),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = rgb
        ln.line.fill.background()
        _set_shape_alpha(ln, alpha)
    # 水平黄金线
    for ratio in (0.382, 0.618):
        ln = slide.shapes.add_shape(
            _SH.RECTANGLE,
            _In(0), _In(H * ratio),
            _In(W), _In(thickness),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = rgb
        ln.line.fill.background()
        _set_shape_alpha(ln, alpha)
