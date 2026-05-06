#!/usr/bin/env python3
"""
pptx_edit.py — v4.2 加载现有 .pptx → AI 改写 → 输出新 .pptx
（对标 ChatPPT Office 插件模式：自然语言修改现有 PPT）

3 种工作模式：
  1. extract  — 解析 pptx 提取所有文字 → 输出 edit-doc JSON（人工改也行）
  2. apply    — 用编辑过的 JSON 把改动写回 pptx（保留所有格式）
  3. rewrite  — 一条龙：extract + Claude API 按 instruction 改写 + apply

保留原 PPT 的：字体 / 颜色 / 字号 / 段落对齐 / 项目符号 / 加粗倾斜
（python-pptx 在 run-level 替换 text，不动其他属性）

用法：
    # 1. 提取文字
    python3 scripts/pptx_edit.py extract /input.pptx --output /tmp/edit.json

    # 2. 人工改完 /tmp/edit.json 后写回
    python3 scripts/pptx_edit.py apply /input.pptx --edits /tmp/edit.json \\
        --output /improved.pptx

    # 3. 一条龙 AI 改写（需 ANTHROPIC_API_KEY）
    python3 scripts/pptx_edit.py rewrite /input.pptx \\
        --instruction "把所有标题改简洁，加 Apple 发布会的克制感" \\
        --output /improved.pptx
"""

from __future__ import annotations
import argparse
import json
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))


def extract_text(pptx_path: Path) -> dict:
    """从 pptx 提取每张 slide 的所有文字 → JSON edit-doc"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("pip install python-pptx")

    prs = Presentation(str(pptx_path))
    slides = []
    for s_idx, slide in enumerate(prs.slides):
        shapes = []
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            paras = []
            for p_idx, para in enumerate(tf.paragraphs):
                runs = [run.text for run in para.runs]
                full = ''.join(runs)
                if not full.strip():
                    continue
                paras.append({
                    'p_idx': p_idx,
                    'text': full,
                    # 记录每个 run 的 text 长度（apply 时用以保留 run-level formatting）
                    'run_lens': [len(r) for r in runs],
                })
            if paras:
                shapes.append({
                    'sh_idx': sh_idx,
                    'shape_name': shape.name,
                    'paragraphs': paras,
                })
        if shapes:
            slides.append({'s_idx': s_idx, 'shapes': shapes})

    return {
        '_pptx': str(pptx_path),
        '_slide_count': len(prs.slides),
        'slides': slides,
    }


def apply_edits(pptx_in: Path, pptx_out: Path, edits: dict) -> int:
    """把 edits JSON 写回 pptx，保留 run-level formatting"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("pip install python-pptx")

    prs = Presentation(str(pptx_in))
    n_changed = 0

    for slide_edit in edits.get('slides', []):
        s_idx = slide_edit['s_idx']
        if s_idx >= len(prs.slides):
            continue
        slide = prs.slides[s_idx]
        shapes_list = list(slide.shapes)

        for shape_edit in slide_edit['shapes']:
            sh_idx = shape_edit['sh_idx']
            if sh_idx >= len(shapes_list):
                continue
            shape = shapes_list[sh_idx]
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            paras_list = list(tf.paragraphs)

            for para_edit in shape_edit['paragraphs']:
                p_idx = para_edit['p_idx']
                if p_idx >= len(paras_list):
                    continue
                para = paras_list[p_idx]
                new_text = para_edit['text']
                runs = list(para.runs)
                if not runs:
                    continue
                # 保留第一个 run 的所有格式，把全部新文本放进第一个 run，清空其他
                old_text = ''.join(r.text for r in runs)
                if new_text == old_text:
                    continue
                # 把全部 new_text 放进第一个 run（保留它的格式）
                runs[0].text = new_text
                for r in runs[1:]:
                    r.text = ''
                n_changed += 1

    prs.save(str(pptx_out))
    return n_changed


def call_claude_rewrite(edits: dict, instruction: str,
                        model: str | None = None) -> dict:
    """让 Claude 按 instruction 改写所有文字，保持 JSON 结构"""
    try:
        import anthropic
    except ImportError:
        raise RuntimeError("pip install anthropic")
    api_key = os.environ.get('ANTHROPIC_API_KEY')
    if not api_key:
        raise RuntimeError("缺 ANTHROPIC_API_KEY")

    from prompt_to_deck import DEFAULT_MODELS
    model = model or os.environ.get('ANTHROPIC_MODEL') or DEFAULT_MODELS['balanced']

    system = """你是 PPT 编辑 AI。输入是 PPT 现有文字的结构化 JSON + 用户的修改指令。
任务：按指令改写每条文字（field "text"），保持 JSON 结构 + 字段 key + 数组顺序完全一致。

【严格规则】
- 不改 _pptx / _slide_count / s_idx / sh_idx / shape_name / p_idx / run_lens 等元字段
- 只改 paragraphs[].text 的文字内容
- 改后文字长度尽量接近原文（避免 PPT 排版崩溃）
- 反 AI Slop：不用"提升 / 优化 / 赋能 / 颗粒度 / 闭环"等公司腔
- 不用 emoji 当 list 标记
- 输出严格 JSON（不要 markdown 包裹）"""

    user = f"""指令：{instruction}

输入 JSON：
```json
{json.dumps(edits, ensure_ascii=False, indent=2)}
```

输出新的完整 JSON。"""

    client = anthropic.Anthropic(api_key=api_key)
    response = client.messages.create(
        model=model,
        max_tokens=8192,
        system=[{"type": "text", "text": system,
                 "cache_control": {"type": "ephemeral"}}],
        messages=[{"role": "user", "content": user}],
    )
    text = response.content[0].text.strip()
    if text.startswith('```'):
        text = text.split('```', 2)[1]
        if text.startswith('json'):
            text = text[4:]
        text = text.strip().rstrip('`').strip()

    new_edits = json.loads(text)

    usage = response.usage
    print(f"  📊 token: input={usage.input_tokens} output={usage.output_tokens}",
          file=sys.stderr)
    if hasattr(usage, 'cache_read_input_tokens') and usage.cache_read_input_tokens:
        print(f"  💾 cache 命中: {usage.cache_read_input_tokens} tokens",
              file=sys.stderr)

    return new_edits


def main():
    parser = argparse.ArgumentParser(description='火一五 PPT v4.2 编辑模式')
    sub = parser.add_subparsers(dest='cmd', required=True)

    p_ext = sub.add_parser('extract', help='提取 pptx 文字 → JSON')
    p_ext.add_argument('pptx')
    p_ext.add_argument('--output', '-o', required=True)

    p_app = sub.add_parser('apply', help='把 JSON 写回 pptx')
    p_app.add_argument('pptx')
    p_app.add_argument('--edits', required=True)
    p_app.add_argument('--output', '-o', required=True)

    p_rw = sub.add_parser('rewrite', help='一条龙 AI 改写')
    p_rw.add_argument('pptx')
    p_rw.add_argument('--instruction', required=True)
    p_rw.add_argument('--output', '-o', required=True)
    p_rw.add_argument('--model', default=None)
    p_rw.add_argument('--save-edits', help='保存中间 JSON（debug 用）')

    args = parser.parse_args()

    if args.cmd == 'extract':
        edits = extract_text(Path(args.pptx))
        Path(args.output).write_text(json.dumps(edits, ensure_ascii=False, indent=2))
        n_p = sum(len(s['paragraphs']) for sl in edits['slides']
                  for s in sl['shapes'])
        print(f"  ✅ {args.output} (extract {edits['_slide_count']} slides, {n_p} paragraphs)",
              file=sys.stderr)
        return

    if args.cmd == 'apply':
        edits = json.loads(Path(args.edits).read_text())
        n = apply_edits(Path(args.pptx), Path(args.output), edits)
        print(f"  ✅ {args.output} ({n} 段文字改动)", file=sys.stderr)
        return

    if args.cmd == 'rewrite':
        print(f"  📤 提取文字...", file=sys.stderr)
        edits = extract_text(Path(args.pptx))
        n_p = sum(len(s['paragraphs']) for sl in edits['slides']
                  for s in sl['shapes'])
        print(f"     {edits['_slide_count']} slides, {n_p} paragraphs", file=sys.stderr)

        print(f"  🤖 Claude 改写: {args.instruction}", file=sys.stderr)
        new_edits = call_claude_rewrite(edits, args.instruction, model=args.model)

        if args.save_edits:
            Path(args.save_edits).write_text(
                json.dumps(new_edits, ensure_ascii=False, indent=2))
            print(f"  💾 中间 JSON: {args.save_edits}", file=sys.stderr)

        print(f"  📥 写回 pptx...", file=sys.stderr)
        n = apply_edits(Path(args.pptx), Path(args.output), new_edits)
        print(f"  ✅ {args.output} ({n} 段改动)", file=sys.stderr)


if __name__ == '__main__':
    main()
