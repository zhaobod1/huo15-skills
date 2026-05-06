#!/usr/bin/env python3
"""
speaker_notes.py — v4.3 自动演讲稿（对标 ChatPPT AI 演讲稿）

读取 .pptx → 提取每张 slide 文字 → 调 Claude 生成 50-150 字讲解词 →
写回 .pptx 备注（notes）层。

5 种语气（说话风格）：
  formal       正式书面，适合董事会 / 学术
  conversational 对话式（默认），适合培训 / 内部分享
  passionate   激情澎湃，适合产品发布 / 路演
  storytelling 讲故事，适合品牌 / 营销
  professor    学者口吻，适合教学

用法：
    python3 scripts/speaker_notes.py /tmp/d.pptx --tone conversational \\
        --output /tmp/d-with-notes.pptx

    # 指定时长（每张 slide 几秒）→ 调整字数
    python3 scripts/speaker_notes.py /tmp/d.pptx --seconds-per-slide 60 \\
        --output /tmp/d-with-notes.pptx
"""

from __future__ import annotations
import argparse
import json
import os
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))


TONE_GUIDE = {
    'formal': '正式书面，第三人称，避免缩写。"本演示将...""数据表明..."',
    'conversational': '对话式自然，第一人称"我们"，可用反问。"这里有个有意思的点..."',
    'passionate': '激情澎湃，短句节奏，感叹。"这就是为什么..."',
    'storytelling': '讲故事开头，"那年...""有一次..."，引发共鸣',
    'professor': '学者口吻，引用数据 / 类比 / 反例。"换个角度看..."',
}


def extract_slide_texts(pptx_path: Path) -> list[dict]:
    """每张 slide 提取所有文字（按 shape 顺序拼接）"""
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = ''.join(r.text for r in para.runs).strip()
                if t:
                    texts.append(t)
        slides.append({
            'idx': i + 1,
            'text': ' / '.join(texts),
        })
    return slides


def generate_notes(slides: list[dict], tone: str, seconds_per_slide: int,
                   model: str | None = None) -> list[str]:
    """调 Claude 一次性生成所有 slide 的 notes"""
    try:
        import anthropic
    except ImportError:
        raise RuntimeError("pip install anthropic")
    api_key = os.environ.get('ANTHROPIC_API_KEY')
    if not api_key:
        raise RuntimeError("缺 ANTHROPIC_API_KEY")
    from prompt_to_deck import DEFAULT_MODELS
    model = model or os.environ.get('ANTHROPIC_MODEL') or DEFAULT_MODELS['balanced']

    # 字数推断：120 字/分钟（中文演讲均速）
    target_chars = max(50, min(300, seconds_per_slide * 2))

    tone_desc = TONE_GUIDE.get(tone, tone)
    system = f"""你是火一五 PPT 演讲稿 AI。给定每张 slide 的文字内容，
为每张写 {target_chars}±20 字的演讲稿（speaker notes），用于讲解员现场照念或参考。

【语气】{tone}：{tone_desc}

【严格规则】
- 输出 JSON 数组：[{{"idx": 1, "notes": "讲解词..."}}, ...]
- notes 字数控制 {target_chars}±20 字
- 不要重复 slide 内已显示的标题文字（要"补充讲解"，不是"复读屏幕"）
- 反 AI Slop：不用"提升 / 优化 / 赋能 / 颗粒度 / 闭环"等公司腔
- 不要"首先 / 其次 / 综上所述"模板腔
- 中文为主，专有名词保留英文"""

    user = f"""为以下 {len(slides)} 张 slide 生成演讲稿：

```json
{json.dumps(slides, ensure_ascii=False, indent=2)}
```

输出严格 JSON 数组（不要 markdown 包裹）。"""

    client = anthropic.Anthropic(api_key=api_key)
    response = client.messages.create(
        model=model,
        max_tokens=8192,
        system=[{"type": "text", "text": system,
                 "cache_control": {"type": "ephemeral"}}],
        messages=[{"role": "user", "content": user}],
    )
    text = response.content[0].text.strip()
    if text.startswith('```'):
        text = text.split('```', 2)[1]
        if text.startswith('json'):
            text = text[4:]
        text = text.strip().rstrip('`').strip()

    notes_data = json.loads(text)
    # 按 idx 排序后取 notes
    notes_data.sort(key=lambda x: x.get('idx', 0))
    notes = [n.get('notes', '') for n in notes_data]

    usage = response.usage
    print(f"  📊 token: input={usage.input_tokens} output={usage.output_tokens}",
          file=sys.stderr)
    return notes


def write_notes(pptx_in: Path, pptx_out: Path, notes: list[str]) -> int:
    """把 notes 写入 .pptx 的备注层"""
    from pptx import Presentation
    prs = Presentation(str(pptx_in))
    n = 0
    for i, slide in enumerate(prs.slides):
        if i >= len(notes):
            break
        text = notes[i].strip()
        if not text:
            continue
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = text
        n += 1
    prs.save(str(pptx_out))
    return n


def main():
    parser = argparse.ArgumentParser(description='火一五 PPT v4.3 自动演讲稿')
    parser.add_argument('pptx', help='输入 .pptx')
    parser.add_argument('--tone', default='conversational',
                        choices=list(TONE_GUIDE.keys()))
    parser.add_argument('--seconds-per-slide', type=int, default=60,
                        help='每张 slide 计划几秒（120 字/分钟换算字数）')
    parser.add_argument('--output', '-o', required=True)
    parser.add_argument('--model', default=None)
    parser.add_argument('--save-notes', help='保存中间 JSON（debug）')
    args = parser.parse_args()

    print(f"  📤 提取 slide 文字...", file=sys.stderr)
    slides = extract_slide_texts(Path(args.pptx))
    print(f"     共 {len(slides)} 张 slide", file=sys.stderr)

    print(f"  🤖 Claude 生成 {args.tone} 语气演讲稿（每张 {args.seconds_per_slide}s）...",
          file=sys.stderr)
    try:
        notes = generate_notes(slides, args.tone, args.seconds_per_slide,
                               model=args.model)
    except Exception as e:
        print(f"  ✗ {e}", file=sys.stderr)
        sys.exit(1)

    if args.save_notes:
        Path(args.save_notes).write_text(
            json.dumps([{'idx': i + 1, 'notes': n}
                       for i, n in enumerate(notes)],
                       ensure_ascii=False, indent=2))
        print(f"  💾 中间 JSON: {args.save_notes}", file=sys.stderr)

    print(f"  📥 写入 .pptx notes 层...", file=sys.stderr)
    n = write_notes(Path(args.pptx), Path(args.output), notes)
    print(f"  ✅ {args.output} ({n} 张 slide 含演讲稿)", file=sys.stderr)


if __name__ == '__main__':
    main()
